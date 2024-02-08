#!/usr/bin/python2

from pptx import Presentation
from pptx.util import Inches, Px, Emu
from os import listdir, makedirs, system
from os.path import isfile, join, exists
import sys

mypath = "images_for_pptx";
lessonname=sys.argv[1]
pdfname=lessonname+".pdf"
if not exists(mypath):
        makedirs(mypath)

system("cd " + mypath + "; ln -s ../" + pdfname + " .; pdftocairo -png " + pdfname + "; rm " + pdfname + "; cd ..")
slide_imgs = []
for f in listdir(mypath):
	filename = join(mypath,f)
	if isfile(filename ) and filename.endswith(".png"):
		slide_imgs.append(filename)
		
slide_imgs.sort()
prs = Presentation()
for s in slide_imgs:
	blank_slidelayout = prs.slidelayouts[6]
	slide = prs.slides.add_slide(blank_slidelayout)
	left = top = Inches(0)
	#pic = slide.shapes.add_picture(s, left, top)
	width = Emu(9144000)
	height = Emu(6858000)
	pic = slide.shapes.add_picture(s, left, top, width, height)


prs.save(lessonname + '.pptx')
system("rm -rf " + mypath)
system("rm -rf *.aux *.loa *.lof *.lot *.log *.nav *.out *.snm *.toc")
